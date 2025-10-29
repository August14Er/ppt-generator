# app.py
import os
import tempfile
from flask import Flask, request, jsonify, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from io import BytesIO
import json

# text extraction
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document

app = Flask(__name__)

# ---------- HELPERS ----------
def save_temp_file(uploaded_file):
    filename = secure_filename(uploaded_file.filename)
    fd, path = tempfile.mkstemp(suffix=os.path.splitext(filename)[1])
    with os.fdopen(fd, 'wb') as tmp:
        tmp.write(uploaded_file.read())
    return path, filename

# ---------- ENDPOINT: extract_text ----------
@app.route('/extract_text', methods=['POST'])
def extract_text():
    if 'file' not in request.files:
        return jsonify({"error":"No file part"}), 400
    f = request.files['file']
    path, filename = save_temp_file(f)
    text = ""
    try:
        if filename.lower().endswith('.pdf'):
            text = extract_pdf_text(path)
        elif filename.lower().endswith('.docx'):
            doc = Document(path)
            text = "\n".join([p.text for p in doc.paragraphs])
        else:
            return jsonify({"error":"Unsupported file type"}), 400
    except Exception as e:
        return jsonify({"error":"Extraction failed","detail": str(e)}), 500
    finally:
        try:
            os.remove(path)
        except:
            pass
    return jsonify({"text": text})

# ---------- ENDPOINT: generate (MODIFICADO) ----------
@app.route('/generate', methods=['POST'])
def generate():
    # 1. Verificar que tengamos los dos archivos
    if 'template_file' not in request.files:
        return jsonify({"error":"No 'template_file' part"}), 400
    if 'slides_data' not in request.form:
        return jsonify({"error":"No 'slides_data' part"}), 400

    template_file = request.files['template_file']
    slides_data_str = request.form['slides_data']

    # 2. Parsear el JSON con los slides
    try:
        slides_data = json.loads(slides_data_str)
        slides = slides_data.get('slides', [])
        if not isinstance(slides, list):
             return jsonify({"error":"'slides' debe ser una lista"}), 400
    except Exception as e:
        return jsonify({"error":"JSON 'slides_data' malformado","detail": str(e)}), 400

    # 3. Cargar la plantilla desde el archivo subido
    try:
        # Usamos BytesIO para leer el archivo en memoria sin guardarlo
        template_stream = BytesIO(template_file.read())
        prs = Presentation(template_stream)
        
        # (El resto de la lógica para añadir slides es igual que antes)
        layout = None
        if len(prs.slide_layouts) > 1:
            layout = prs.slide_layouts[1]
        else:
            layout = prs.slide_layouts[0]

        for s in slides:
            title = s.get('title', '')
            body = s.get('body', '')
            slide = prs.slides.add_slide(layout)
            try:
                if slide.shapes.title:
                    slide.shapes.title.text = title
            except Exception: pass
            
            filled = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    try:
                        if shape is slide.shapes.title:
                            continue
                    except Exception: pass
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = body
                    filled = True
                    break
            if not filled:
                from pptx.util import Inches, Pt
                left = Inches(1); top = Inches(1.5); width = Inches(8); height = Inches(4.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = body

        # 4. Guardar y devolver el PPT final
        out = BytesIO()
        prs.save(out)
        out.seek(0)
        return send_file(out,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True,
                         download_name='presentation_generated.pptx')
                         
    except Exception as e:
        return jsonify({"error":"Fallo en la generación del PPT","detail": str(e)}), 500

# ---------- Root for health check ----------
@app.route('/')
def home():
    return jsonify({"message":"PPT Generator Flask App is running!"})

if __name__ == '__main__':
    # Render requires using the PORT env var
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)

